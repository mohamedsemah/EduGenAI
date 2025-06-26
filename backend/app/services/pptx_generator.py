from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from app.models.lesson import LessonContent
import os


def create_presentation(lesson_content: LessonContent, output_path: str):
    """Create a PowerPoint presentation based on lesson content"""

    # Create a new presentation
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content

    # Create title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = lesson_content.title
    subtitle.text = f"Grade Level: {lesson_content.grade_level}\nDuration: {lesson_content.duration}\n\nBased on Universal Design for Learning (UDL) Principles"

    # Style the title slide
    try:
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # Dark blue

        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)  # Gray
    except:
        pass  # If styling fails, continue without it

    # Add note to title slide with accessibility info
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = (
        f"This presentation follows Universal Design for Learning principles and includes "
        f"the following accessibility features:\n\n"
        f"Visual: {lesson_content.accessibility_features.get('visual', 'Standard visual formatting')}\n"
        f"Auditory: {lesson_content.accessibility_features.get('auditory', 'Clear verbal instructions')}\n"
        f"Cognitive: {lesson_content.accessibility_features.get('cognitive', 'Simple structure and language')}\n"
        f"Physical: {lesson_content.accessibility_features.get('physical', 'Standard navigation')}\n"
        f"Language: {lesson_content.accessibility_features.get('language', 'Clear language support')}\n\n"
        f"Lesson Overview: {lesson_content.overview}"
    )

    # Create slide for each content slide
    for slide_content in lesson_content.slides:
        slide = prs.slides.add_slide(content_slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_content.title

        # Style the title
        try:
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        except:
            pass

        # Set content
        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()  # Clear existing content

            # Add the main content
            p = text_frame.paragraphs[0]
            p.text = slide_content.content

            # Style the content
            try:
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(52, 73, 94)
                text_frame.margin_left = Inches(0.5)
                text_frame.margin_right = Inches(0.5)
                text_frame.margin_top = Inches(0.3)
                text_frame.margin_bottom = Inches(0.3)
            except:
                pass

        # Add a text box with accessibility information at the bottom
        try:
            accessibility_info = []
            for key, value in slide_content.accessibility_features.items():
                accessibility_info.append(f"{key.replace('_', ' ').title()}: {value}")

            if accessibility_info:
                # Add small text box at bottom for accessibility features
                left = Inches(0.5)
                top = Inches(7)
                width = Inches(9)
                height = Inches(0.8)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = "Accessibility Features: " + " | ".join(accessibility_info[:3])  # Limit to first 3

                # Style accessibility text
                p = text_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(149, 165, 166)
                p.font.italic = True
        except:
            pass  # If accessibility text box fails, continue without it

        # Add notes with instructions and accessibility features
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        # Build comprehensive notes
        notes_content = []

        # Add presenter notes
        if slide_content.notes:
            notes_content.append("PRESENTER NOTES:")
            notes_content.append(slide_content.notes)
            notes_content.append("")

        # Add image description
        if slide_content.image_prompt:
            notes_content.append("VISUAL DESCRIPTION:")
            notes_content.append(f"Recommended image: {slide_content.image_prompt}")
            notes_content.append("")

        # Add accessibility features
        if slide_content.accessibility_features:
            notes_content.append("ACCESSIBILITY FEATURES:")
            for feature, description in slide_content.accessibility_features.items():
                notes_content.append(f"• {feature.replace('_', ' ').title()}: {description}")
            notes_content.append("")

        # Add UDL principles reminder
        notes_content.append("UDL PRINCIPLES APPLIED:")
        notes_content.append("• Multiple means of representation (visual, auditory, text)")
        notes_content.append("• Multiple means of engagement (choice, relevance, motivation)")
        notes_content.append("• Multiple means of action/expression (flexible demonstration of learning)")

        text_frame.text = "\n".join(notes_content)

    # Add a final resources slide
    resource_slide = prs.slides.add_slide(content_slide_layout)
    title = resource_slide.shapes.title
    title.text = "Additional Resources & Support"

    if len(resource_slide.placeholders) > 1:
        content = resource_slide.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()

        resource_content = f"""📚 Materials Used in This Lesson:
{chr(10).join([f'• {material}' for material in lesson_content.materials])}

🎯 Assessment Strategies:
• {lesson_content.assessment}

🔄 Next Steps:
• {lesson_content.conclusion}

♿ Accessibility Support:
• All slides include alt text and high contrast
• Content available in multiple formats
• Flexible pacing and participation options
• Assistive technology compatible

📞 For additional support or accommodations, please contact your instructor."""

        p = text_frame.paragraphs[0]
        p.text = resource_content

        try:
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(52, 73, 94)
        except:
            pass

    # Add final notes
    notes_slide = resource_slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = f"""This lesson plan was generated using Universal Design for Learning principles to ensure accessibility and engagement for all learners.

Key UDL Elements Included:
- Multiple means of representation
- Multiple means of engagement  
- Multiple means of action and expression

The lesson is designed for {lesson_content.grade_level} students and should take approximately {lesson_content.duration} to complete.

All materials and activities can be adapted further based on specific student needs and classroom contexts."""

    # Save the presentation
    try:
        prs.save(output_path)
        print(f"Presentation saved successfully to {output_path}")
        return output_path
    except Exception as e:
        print(f"Error saving presentation: {e}")
        raise e